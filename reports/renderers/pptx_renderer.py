from __future__ import annotations

from io import BytesIO
from typing import Any, Dict, List, Optional

from pptx import Presentation

LAYOUTS = {
    "title": 0,        # Title Slide
    "bullets": 1,      # Title and Content
    "two_columns": 3,  # Two Content
    "sources": 1,      # Title and Content
}


def render_pptx_from_spec(spec: Dict[str, Any], *, debug_layouts: bool = False, use_letterhead: bool = False) -> bytes:
    prs = None
    if use_letterhead: # si se va a usar plantilla
        import os
        from django.conf import settings
        template_path = os.path.join(settings.BASE_DIR, "reports", "templates", "plantilla_secihti.pptx")
        if os.path.exists(template_path):
            prs = Presentation(template_path)
            # Remove all existing sample slides from the template
            xml_slides = prs.slides._sldIdLst  
            for s in list(xml_slides):
                xml_slides.remove(s)
        else:
            print(f"[REPORT] ADVERTENCIA: Plantilla {template_path} no encontrada.")
            prs = Presentation()
    else:
        prs = Presentation()

    if debug_layouts:
        _debug_print_layouts(prs)

    slides = spec.get("slides", [])
    if not isinstance(slides, list):
        raise ValueError("spec['slides'] debe ser una lista")

    for slide_spec in slides:
        layout_name = slide_spec.get("layout", "bullets")
        layout_idx = LAYOUTS.get(layout_name, 1)   #1, se pasa por default, , significa  titulo y contendido

        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        # Título
        if slide.shapes.title and slide_spec.get("title"):
            slide.shapes.title.text = str(slide_spec["title"])

        # Ruteo por layout
        if layout_name == "title":
            _set_subtitle(slide, slide_spec.get("subtitle"))

        elif layout_name == "bullets":
            body = _get_body_placeholder(slide)
            _fill_bullets(body, slide_spec.get("bullets", []))
            _set_notes(slide, slide_spec.get("notes"))

        elif layout_name == "two_columns":
            placeholders = list(slide.placeholders)
            left_ph = _find_placeholder_by_idx(placeholders, 1)
            right_ph = _find_placeholder_by_idx(placeholders, 2)

            if left_ph:
                _fill_column(left_ph, slide_spec.get("left", {}))
            if right_ph:
                _fill_column(right_ph, slide_spec.get("right", {}))

            _set_notes(slide, slide_spec.get("notes"))

        elif layout_name == "sources":
            body = _get_body_placeholder(slide)
            sources = slide_spec.get("sources", [])
            lines = [f'{s.get("name","")}: {s.get("detail","")}' for s in sources]
            _fill_bullets(body, lines)

        else:
            # fallback: trata como bullets
            body = _get_body_placeholder(slide)
            _fill_bullets(body, slide_spec.get("bullets", []))
            _set_notes(slide, slide_spec.get("notes"))

    out = BytesIO()
    prs.save(out)
    return out.getvalue()


def _debug_print_layouts(prs: Presentation) -> None:
    print("== PPTX Layouts disponibles ==")
    for i, layout in enumerate(prs.slide_layouts):
        name = getattr(layout, "name", None) or f"layout_{i}"
        print(i, "-", name)


def _set_subtitle(slide, subtitle: Optional[str]) -> None:
    if not subtitle:
        return
    # En "Title Slide", el subtítulo suele ser placeholder idx=1
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            ph.text = str(subtitle)
            return


def _get_body_placeholder(slide):
    # BODY suele existir en layouts "Title and Content"
    for ph in slide.placeholders:
        # BODY=2
        if ph.placeholder_format.type == 2:
            return ph
    # fallback: primer placeholder que no sea el título
    for ph in slide.placeholders:
        if slide.shapes.title and ph != slide.shapes.title:
            return ph
    return None


def _fill_bullets(body_placeholder, bullets: List[str]) -> None:
    if body_placeholder is None:
        return

    if bullets is None:
        bullets = []
    if not isinstance(bullets, list):
        raise ValueError("'bullets' debe ser una lista de strings")

    tf = body_placeholder.text_frame
    
    # Safe clear to retain master placeholder formatting
    while len(tf.paragraphs) > 1:
        p = tf.paragraphs[-1]
        p._element.getparent().remove(p._element)
        
    if len(tf.paragraphs) > 0:
        tf.paragraphs[0].text = ""

    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = str(text)
        p.level = 0


def _find_placeholder_by_idx(placeholders, idx: int):
    for p in placeholders:
        if p.placeholder_format.idx == idx:
            return p
    return None


def _fill_column(placeholder, data: Dict[str, Any]) -> None:
    tf = placeholder.text_frame
    
    while len(tf.paragraphs) > 1:
        p = tf.paragraphs[-1]
        p._element.getparent().remove(p._element)
        
    if len(tf.paragraphs) > 0:
        tf.paragraphs[0].text = ""

    heading = data.get("heading")
    bullets = data.get("bullets", [])

    if heading:
        if len(tf.paragraphs) == 0:
            tf.add_paragraph()
        tf.paragraphs[0].text = str(heading)

    if bullets and not isinstance(bullets, list):
        raise ValueError("column['bullets'] debe ser una lista")

    start_idx = 1 if heading else 0
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if (i == 0 and not heading) else tf.add_paragraph()
        p.text = str(b)
        p.level = 1


def _set_notes(slide, notes: Optional[str]) -> None:
    if not notes:
        return
    slide.notes_slide.notes_text_frame.text = str(notes)
